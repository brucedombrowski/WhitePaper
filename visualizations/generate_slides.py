#!/usr/bin/env python3
"""
Generate a fun, educational, inspiring slide deck for Friday team meeting.
Topic: Git-based version control workflow for non-technical teams.

Uses python-pptx to create PowerPoint with embedded charts and cartoon-style graphics.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from pathlib import Path

OUTPUT_DIR = Path(__file__).parent
CHARTS_DIR = OUTPUT_DIR

# Color palette (navy/blue theme matching the project)
NAVY = RGBColor(0x0F, 0x27, 0x44)
BLUE = RGBColor(0x1E, 0x49, 0x76)
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xC4)
ACCENT = RGBColor(0x8E, 0xAC, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
RED = RGBColor(0xD6, 0x27, 0x28)
ORANGE = RGBColor(0xFF, 0x7F, 0x0E)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_background(slide, color=NAVY):
    """Fill slide background with color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, items, start_top=2.2, font_size=22, color=WHITE,
                     emoji_list=None, left=1.5):
    """Add bullet points with optional emojis."""
    for i, item in enumerate(items):
        prefix = emoji_list[i] if emoji_list and i < len(emoji_list) else "\u2022"
        add_textbox(slide, left, start_top + i * 0.65, 10, 0.6,
                    f"{prefix}  {item}", font_size=font_size, color=color)


def add_cartoon_shape(slide, shape_type, left, top, width, height,
                      fill_color, text="", font_size=36):
    """Add a cartoon-style shape with text."""
    shape = slide.shapes.add_shape(shape_type,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = WHITE
    shape.line.width = Pt(3)
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_arrow(slide, left, top, width, height, color=LIGHT_BLUE):
    """Add a right-pointing arrow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_image_safe(slide, path, left, top, width=None, height=None):
    """Add image if it exists."""
    if os.path.exists(path):
        kwargs = {}
        if width:
            kwargs['width'] = Inches(width)
        if height:
            kwargs['height'] = Inches(height)
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)
        return True
    return False


# ============================================================================
# SLIDE 1: Title
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, NAVY)

add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1, 0.5, 11.3, 2.5, BLUE)
add_textbox(slide, 1.5, 0.7, 10, 1.2,
            "Version Control for Everyone",
            font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 1.8, 10, 0.8,
            "How Git + Your Browser = Superpowers",
            font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Fun cartoon icons row
shapes_data = [
    (MSO_SHAPE.OVAL, 2, 3.8, 1.5, 1.5, LIGHT_BLUE, "\U0001F4BB"),      # laptop
    (MSO_SHAPE.OVAL, 4.2, 3.8, 1.5, 1.5, GREEN, "\U0001F500"),          # arrows
    (MSO_SHAPE.OVAL, 6.4, 3.8, 1.5, 1.5, ORANGE, "\U0001F310"),         # globe
    (MSO_SHAPE.OVAL, 8.6, 3.8, 1.5, 1.5, RGBColor(0x94, 0x67, 0xBD), "\u2705"), # check
]
for shape_type, l, t, w, h, c, txt in shapes_data:
    add_cartoon_shape(slide, shape_type, l, t, w, h, c, txt, font_size=40)

# Labels under circles
labels = ["You Edit", "Git Tracks", "Team Reviews", "Approved!"]
for i, label in enumerate(labels):
    add_textbox(slide, 2 + i * 2.2, 5.4, 1.5, 0.5, label,
                font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Arrows between circles
for i in range(3):
    add_arrow(slide, 3.55 + i * 2.2, 4.35, 0.6, 0.4, ACCENT)

add_textbox(slide, 1.5, 6.3, 10, 0.5,
            "Friday Training  |  Bruce Dombrowski",
            font_size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 6.8, 10, 0.5,
            "No CLI required. Just your browser.",
            font_size=16, color=RGBColor(0x66, 0x88, 0xAA), alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 2: The Problem
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.5, 0.3, 12, 1,
            "\U0001F62B  The Problem: File Chaos",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Left side: chaotic filenames
add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, 1.5, 5.5, 5.2, RGBColor(0x8B, 0x00, 0x00))
add_textbox(slide, 1.2, 1.7, 5, 0.5,
            "Sound familiar?", font_size=24, color=YELLOW, bold=True)

chaos_files = [
    "budget_FINAL.xlsx",
    "budget_FINAL_v2.xlsx",
    "budget_FINAL_v2_REVISED.xlsx",
    "budget_FINAL_v2_REVISED_FIXED.xlsx",
    "budget_ACTUAL_FINAL_USE_THIS.xlsx",
    "budget_ACTUAL_FINAL_USE_THIS(1).xlsx",
    "Copy of budget_DONT_DELETE.xlsx",
]
for i, f in enumerate(chaos_files):
    add_textbox(slide, 1.5, 2.4 + i * 0.55, 5, 0.5,
                f"\U0001F4C4  {f}", font_size=16, color=WHITE)

# Right side: the questions
add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.8, 1.5, 5.5, 5.2, BLUE)
add_textbox(slide, 7.2, 1.7, 5, 0.5,
            "The real questions:", font_size=24, color=YELLOW, bold=True)

questions = [
    ("\U0001F914", "Which version is current?"),
    ("\U0001F50D", "What changed between versions?"),
    ("\U0001F464", "Who made this change?"),
    ("\U0001F4C5", "When was it changed?"),
    ("\u2753", "Why was it changed?"),
    ("\U0001F512", "Can we prove the audit trail?"),
    ("\U0001F91D", "Did the team approve this?"),
]
for i, (emoji, q) in enumerate(questions):
    add_textbox(slide, 7.5, 2.4 + i * 0.55, 5, 0.5,
                f"{emoji}  {q}", font_size=18, color=WHITE)


# ============================================================================
# SLIDE 3: The Solution
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.5, 0.3, 12, 1,
            "\U0001F389  The Solution: Git + Browser",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 1.2, 12, 0.6,
            "Git tracks every change. Your browser shows it. No software to install.",
            font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Three big benefit boxes
benefits = [
    ("\U0001F4DC", "Complete\nHistory", "Every version saved\nautomatically.\nNothing lost. Ever.", GREEN),
    ("\U0001F50E", "See What\nChanged", "Line-by-line diffs.\nGreen = added.\nRed = removed.", LIGHT_BLUE),
    ("\U0001F46B", "Team\nReview", "Everyone reviews\nbefore it goes live.\nBuilt-in approval.", ORANGE),
]
for i, (emoji, title, desc, color) in enumerate(benefits):
    x = 1 + i * 4
    add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.2, 3.5, 4.5, color)
    add_textbox(slide, x + 0.2, 2.3, 3.1, 1, emoji, font_size=60,
                color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 3.3, 3.1, 1, title, font_size=28,
                color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.3, 4.5, 2.9, 2, desc, font_size=18,
                color=WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 4: How It Works (5 steps)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.5, 0.3, 12, 1,
            "\U0001F3AF  How It Works: 5 Easy Steps",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

steps = [
    ("\U0001F4DD", "Edit", "Edit your file\n(Excel, CSV,\nwhatever)", GREEN),
    ("\U0001F4BE", "Save", "Save like\nyou always do", LIGHT_BLUE),
    ("\U0001F4E4", "Commit", "Click 'Commit'\n(one click saves\nthe version)", ORANGE),
    ("\U0001F4AC", "Review", "Team reviews\nin their browser\n(comments + diffs)", RGBColor(0x94, 0x67, 0xBD)),
    ("\u2705", "Approve", "Approved = merged\nto the official\nversion", GREEN),
]
for i, (emoji, title, desc, color) in enumerate(steps):
    x = 0.5 + i * 2.5
    add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.8, 2.2, 4.5, color)
    add_textbox(slide, x + 0.1, 1.9, 2, 0.8, emoji, font_size=50,
                color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 2.8, 2, 0.6, title, font_size=26,
                color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 3.6, 1.8, 2.5, desc, font_size=16,
                color=WHITE, alignment=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < 4:
        add_arrow(slide, x + 2.25, 3.8, 0.35, 0.35, ACCENT)

add_textbox(slide, 0.5, 6.6, 12, 0.6,
            "\U0001F4A1  The whole team does this in their browser. No command line. No special software.",
            font_size=20, color=YELLOW, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 5: What the Browser Shows You
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.5, 0.3, 12, 1,
            "\U0001F4BB  What You See in Your Browser",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Simulate a diff view
add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.5, 1.5, 10.3, 5.3,
                  RGBColor(0x1E, 0x1E, 0x2E))

add_textbox(slide, 2, 1.7, 9, 0.5,
            "\U0001F4C1  budget.csv  — Changed 3 lines",
            font_size=20, color=ACCENT, bold=True)

# Simulated diff lines
diff_lines = [
    ("  ", "Department, Q1 Budget, Q2 Budget, Status", RGBColor(0xA0, 0xA0, 0xA0)),
    ("  ", "Engineering, $150000, $155000, Active", RGBColor(0xA0, 0xA0, 0xA0)),
    ("- ", "Marketing,  $80000,  $80000,  Active", RGBColor(0xFF, 0x60, 0x60)),
    ("+ ", "Marketing,  $80000,  $95000,  Active    \u2190 budget increased!", RGBColor(0x60, 0xFF, 0x60)),
    ("  ", "Operations, $200000, $200000, Active", RGBColor(0xA0, 0xA0, 0xA0)),
    ("- ", "Sales,      $120000, $120000, Active", RGBColor(0xFF, 0x60, 0x60)),
    ("+ ", "Sales,      $120000, $135000, Active    \u2190 budget increased!", RGBColor(0x60, 0xFF, 0x60)),
    ("  ", "HR,         $90000,  $90000,  Active", RGBColor(0xA0, 0xA0, 0xA0)),
    ("- ", "IT,         $175000, $175000, Pending", RGBColor(0xFF, 0x60, 0x60)),
    ("+ ", "IT,         $175000, $175000, Active    \u2190 status changed!", RGBColor(0x60, 0xFF, 0x60)),
]
for i, (prefix, line, color) in enumerate(diff_lines):
    # Background highlight for changed lines
    if prefix == "- ":
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(1.8), Inches(2.35 + i * 0.38),
                                     Inches(9.8), Inches(0.36))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x50, 0x10, 0x10)
        bg.line.fill.background()
    elif prefix == "+ ":
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(1.8), Inches(2.35 + i * 0.38),
                                     Inches(9.8), Inches(0.36))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x10, 0x40, 0x10)
        bg.line.fill.background()

    add_textbox(slide, 2, 2.35 + i * 0.38, 9.5, 0.4,
                f"{prefix}{line}", font_size=14, color=color,
                font_name='Courier New')

add_textbox(slide, 1.5, 6.5, 10, 0.5,
            "\U0001F449  Instantly see exactly what changed, who changed it, and when.",
            font_size=20, color=YELLOW, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 6: Excel Workflow
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.5, 0.3, 12, 1,
            "\U0001F4CA  But What About My Excel Files?",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 1.2, 12, 0.6,
            "Excel is binary (can't see diffs). Solution: auto-convert to CSV!",
            font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Before/After
add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, 2.2, 5.2, 4.5, RGBColor(0x8B, 0x00, 0x00))
add_textbox(slide, 1, 2.3, 5, 0.5,
            "\U0001F6AB  Before: Can't See Changes", font_size=22, color=WHITE, bold=True)
before_items = [
    "\U0001F4C4  budget.xlsx uploaded",
    "\U0001F4AC  \"Binary file changed\"",
    "\U0001F937  No idea what's different",
    "\U0001F612  \"Just trust me, I updated it\"",
    "\U0001F4A2  Audit nightmare",
]
for i, item in enumerate(before_items):
    add_textbox(slide, 1.3, 3.0 + i * 0.6, 4.5, 0.5, item, font_size=17, color=WHITE)

add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7, 2.2, 5.5, 4.5, RGBColor(0x00, 0x6B, 0x00))
add_textbox(slide, 7.2, 2.3, 5, 0.5,
            "\u2705  After: Full Visibility", font_size=22, color=WHITE, bold=True)
after_items = [
    "\U0001F4C4  budget.xlsx + budget.csv",
    "\U0001F50E  CSV shows every change",
    "\U0001F7E2  Green = added, \U0001F534 Red = removed",
    "\U0001F91D  Team approves in browser",
    "\U0001F3C6  Audit-ready automatically",
]
for i, item in enumerate(after_items):
    add_textbox(slide, 7.3, 3.0 + i * 0.6, 5, 0.5, item, font_size=17, color=WHITE)

add_arrow(slide, 6.1, 4.0, 0.8, 0.5, YELLOW)

add_textbox(slide, 0.5, 6.8, 12, 0.5,
            "\U0001F4A1  A pre-commit hook auto-converts Excel to CSV. You never think about it.",
            font_size=18, color=YELLOW, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 7: Branch & Review Workflow
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.5, 0.3, 12, 1,
            "\U0001F333  Branches: Safe Experimentation",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 1.2, 12, 0.6,
            "Make changes in a branch. Team reviews. Merge when approved.",
            font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Visual branch diagram using shapes
# Main branch line
main_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1), Inches(3.5),
                                    Inches(11), Inches(0.15))
main_line.fill.solid()
main_line.fill.fore_color.rgb = GREEN
main_line.line.fill.background()

add_textbox(slide, 0.2, 3.3, 0.8, 0.4, "main", font_size=16, color=GREEN, bold=True)

# Branch line (curves up)
branch_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(3.5), Inches(2.5),
                                      Inches(5), Inches(0.12))
branch_line.fill.solid()
branch_line.fill.fore_color.rgb = ORANGE
branch_line.line.fill.background()

add_textbox(slide, 2.5, 2.15, 2, 0.4, "your-changes", font_size=16, color=ORANGE, bold=True)

# Commit dots on main
for x in [1.5, 3.0, 9.0, 10.5]:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x), Inches(3.35),
                                  Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GREEN
    dot.line.color.rgb = WHITE
    dot.line.width = Pt(2)

# Commit dots on branch
for x in [4.0, 5.5, 7.0]:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x), Inches(2.35),
                                  Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ORANGE
    dot.line.color.rgb = WHITE
    dot.line.width = Pt(2)

# Labels for key points
labels_data = [
    (3.5, 4.2, "\U0001F331 Branch off", ORANGE),
    (5.0, 1.8, "\U0001F4DD Edit, edit, edit", ORANGE),
    (7.5, 2.0, "\U0001F4AC Team reviews\n(Merge Request)", RGBColor(0x94, 0x67, 0xBD)),
    (8.5, 4.2, "\u2705 Merge!", GREEN),
]
for x, y, text, color in labels_data:
    add_textbox(slide, x, y, 2.5, 0.8, text, font_size=15, color=color, alignment=PP_ALIGN.CENTER)

# Bottom callout
add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.5, 5.3, 10, 1.8, BLUE)
add_textbox(slide, 2, 5.4, 9, 0.5,
            "\U0001F6E1\uFE0F  Safety Guardrails (set once, enforced forever):",
            font_size=20, color=YELLOW, bold=True)
guardrails = [
    "\U0001F512  No one can push directly to main — all changes go through review",
    "\U0001F46B  Required approvers must sign off before merge",
    "\U0001F504  If you change your code after approval, approvals reset automatically",
]
for i, g in enumerate(guardrails):
    add_textbox(slide, 2.2, 6.0 + i * 0.35, 9, 0.35, g, font_size=15, color=WHITE)


# ============================================================================
# SLIDE 8: Why This Matters (Compliance)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.5, 0.3, 12, 1,
            "\U0001F3DB\uFE0F  Why This Matters: Built-in Compliance",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 1.2, 12, 0.6,
            "Everything the auditor needs is automatic. No extra paperwork.",
            font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Compliance mapping boxes
compliance = [
    ("\U0001F4DD", "Change Control\n(NIST CM-3)", "Every change tracked\nwith who, what, when,\nand why", GREEN),
    ("\U0001F464", "Separation of\nDuties (AC-5)", "Author cannot approve\ntheir own changes.\nSystematically enforced.", LIGHT_BLUE),
    ("\U0001F4CB", "Audit Trail\n(AU-3)", "Complete, timestamped,\ntamper-evident history.\nNo gaps possible.", ORANGE),
    ("\U0001F512", "Access Control\n(AC-6)", "Least privilege.\nView-only for reviewers.\nRole-based access.", RGBColor(0x94, 0x67, 0xBD)),
]
for i, (emoji, title, desc, color) in enumerate(compliance):
    x = 0.5 + i * 3.2
    add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.2, 2.8, 4.5, color)
    add_textbox(slide, x + 0.1, 2.3, 2.6, 0.8, emoji, font_size=50,
                color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 3.2, 2.6, 1, title, font_size=20,
                color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 4.4, 2.4, 2, desc, font_size=16,
                color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 6.8, 12, 0.5,
            "\U0001F4A1  The auditor opens a URL. Everything is there. Meeting over in 5 minutes.",
            font_size=20, color=YELLOW, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 9: Chart - Repo Ecosystem (embed our chart)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.5, 0.3, 12, 1,
            "\U0001F4CA  What This Looks Like in Practice",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

chart_path = CHARTS_DIR / 'repo_comparison.png'
if add_image_safe(slide, chart_path, 0.8, 1.5, width=11.5):
    add_textbox(slide, 0.5, 6.0, 12, 0.5,
                "636 commits across 7 repositories in 3 weeks. All tracked. All auditable.",
                font_size=20, color=YELLOW, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 6.5, 12, 0.5,
                "One engineer + AI agents = this output.",
                font_size=18, color=ACCENT, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 10: Chart - Ecosystem Timeline
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.5, 0.3, 12, 1,
            "\U0001F4C5  When Things Happened",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

chart_path = CHARTS_DIR / 'ecosystem_timeline.png'
if add_image_safe(slide, chart_path, 0.8, 1.5, width=11.5):
    add_textbox(slide, 0.5, 5.8, 12, 0.5,
                "Each bar = one project's active development window. All running concurrently.",
                font_size=20, color=YELLOW, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 11: Chart - Daily Activity
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.5, 0.3, 12, 1,
            "\U0001F4C8  Daily Commit Activity",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

chart_path = CHARTS_DIR / 'daily_activity.png'
if add_image_safe(slide, chart_path, 1, 1.5, width=11):
    add_textbox(slide, 0.5, 5.6, 12, 0.5,
                "Every bar = a day of tracked, auditable work. Colors = different projects.",
                font_size=20, color=YELLOW, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 12: Gource Snapshot
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.5, 0.3, 12, 0.8,
            "\U0001F30C  The Whole Ecosystem, Visualized",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

chart_path = CHARTS_DIR / 'gource-snapshot.png'
if add_image_safe(slide, chart_path, 1.5, 1.2, width=10):
    add_textbox(slide, 0.5, 6.0, 12, 0.5,
                "636 commits. 7 repos. 34,000+ lines of code. 3 weeks. One engineer.",
                font_size=22, color=YELLOW, alignment=PP_ALIGN.CENTER, bold=True)
    add_textbox(slide, 0.5, 6.5, 12, 0.5,
                "(We have a 40-second animated video of this too \U0001F3AC)",
                font_size=18, color=ACCENT, alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 13: Getting Started
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_textbox(slide, 0.5, 0.3, 12, 1,
            "\U0001F680  Getting Started: Your First 5 Minutes",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

steps_start = [
    ("\U0001F4E7", "You'll get an email invitation to the GitLab/GitHub project", GREEN),
    ("\U0001F310", "Click the link — opens in your browser", LIGHT_BLUE),
    ("\U0001F4C2", "Browse files, click any file to see its contents", ORANGE),
    ("\U0001F553", "Click \"History\" to see every version ever saved", RGBColor(0x94, 0x67, 0xBD)),
    ("\U0001F50D", "Click any version to see exactly what changed (the diff)", GREEN),
    ("\U0001F4AC", "Leave a comment if you have questions or concerns", LIGHT_BLUE),
    ("\u2705", "Click \"Approve\" when you're satisfied with the changes", ORANGE),
]
for i, (emoji, text, color) in enumerate(steps_start):
    # Number circle
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(1), Inches(1.5 + i * 0.75),
                                        Inches(0.55), Inches(0.55))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = color
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, 1.8, 1.52 + i * 0.75, 10, 0.5,
                f"{emoji}  {text}", font_size=20, color=WHITE)


# ============================================================================
# SLIDE 14: Closing
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.5, 1, 10.3, 3, BLUE)
add_textbox(slide, 2, 1.2, 9, 1.2,
            "\U0001F31F  You Already Know How to Do This",
            font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 2, 2.5, 9, 1,
            "If you can browse the web, you can review changes.\nIf you can leave a comment, you can approve a release.\nIf you can click a link, you have an audit trail.",
            font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Bottom section
add_textbox(slide, 0.5, 4.5, 12, 0.8,
            "What used to require:",
            font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)

# Before vs After
add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.5, 5.3, 4.5, 1.5, RGBColor(0x8B, 0x00, 0x00))
add_textbox(slide, 1.7, 5.4, 4, 1.3,
            "\u274C  Spreadsheets of spreadsheets\n\u274C  Email chains of approvals\n\u274C  \"Who has the latest version?\"",
            font_size=17, color=WHITE)

add_arrow(slide, 6.2, 5.7, 0.8, 0.5, YELLOW)

add_cartoon_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.3, 5.3, 4.5, 1.5, RGBColor(0x00, 0x6B, 0x00))
add_textbox(slide, 7.5, 5.4, 4, 1.3,
            "\u2705  One URL\n\u2705  Complete history\n\u2705  Built-in compliance",
            font_size=17, color=WHITE)

add_textbox(slide, 0.5, 7.0, 12, 0.5,
            "Questions?  \U0001F64B",
            font_size=28, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================================
# Save
# ============================================================================
output_path = OUTPUT_DIR / 'git-workflow-training.pptx'
prs.save(str(output_path))
print(f'Saved: {output_path}')
print(f'Slides: {len(prs.slides)}')
print(f'Size: {os.path.getsize(output_path) / 1024:.0f} KB')
